"""presentation_write çok-slaytlı deste türetimi — "5 sayfalık sunum" sözleşmesi.

Ekran görüntüsündeki hata sınıfı: slides verilmeyince başlık + tek içerik
slaytı basılıyor, istenen sayfa sayısı yok sayılıyor ve içerik kopyalanmış
görünüyordu. Türetici, araştırma özetini (numaralı bulgular + Kaynaklar
kuyruğu) istenen sayıda slayta dağıtmalı.
"""

from __future__ import annotations

from pathlib import Path

from actions.presentation_write import (
    _derive_slide_specs,
    _requested_slide_count,
    presentation_write,
)

_RESEARCH_SUMMARY = (
    "1. Kuantum interneti büyük adım: dolanıklık ilk kez şehirlerarası taşındı.\n"
    "2. Kuantum tekrarlayıcılar menzil sorununu çözmeye yaklaşıyor.\n"
    "3. Standartlaşma çalışmaları ITU bünyesinde başladı.\n"
    "4. Türkiye'de ilk kuantum ağı pilotu duyuruldu.\n"
    "Kaynaklar: - example.com/kuantum - example.org/haber"
)


def test_requested_slide_count_parses_turkish_forms() -> None:
    assert _requested_slide_count("5 sayfalık sunum hazırla") == 5
    assert _requested_slide_count("3 slaytlık deste") == 3
    assert _requested_slide_count("sunum hazırla") == 0
    # Aşırı istekler sınırlanır.
    assert _requested_slide_count("99 sayfalık sunum") == 12


def test_derive_slide_specs_distributes_findings() -> None:
    specs = _derive_slide_specs(_RESEARCH_SUMMARY, "Kuantum", target=5)

    # başlık slaytı yazıcı tarafında ayrı: türetici 5-1=4 slot → 3 içerik + Kaynaklar
    assert len(specs) == 4
    assert specs[-1]["title"] == "Kaynaklar"
    all_bullets = [b for spec in specs[:-1] for b in spec["bullets"]]
    assert len(all_bullets) == 4  # dört bulgu da dağıtıldı
    assert all(spec["bullets"] for spec in specs)  # boş slayt yok


def test_presentation_write_honors_page_target(tmp_path: Path) -> None:
    outcome = presentation_write(
        prompt="kuantum interneti hakkında 5 sayfalık sunum hazırla",
        output_path=str(tmp_path / "kuantum.pptx"),
        title="kuantum interneti",
        source_context=_RESEARCH_SUMMARY,
        overwrite=True,
    )
    assert "PPTX oluşturuldu" in outcome["text"]

    from pptx import Presentation

    deck = Presentation(str(tmp_path / "kuantum.pptx"))
    slides = list(deck.slides)
    assert len(slides) == 5
    # İçerik slaytları birbirinin kopyası olmamalı.
    def slide_text(slide):
        return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    contents = [slide_text(s) for s in slides[1:-1]]
    assert len(set(contents)) == len(contents)
    assert all(content.count("•") <= 5 for content in contents)
    assert deck.slide_width > deck.slide_height

    body_sizes = []
    for slide in slides[1:-1]:
        for shape in slide.shapes:
            if not shape.has_text_frame or "•" not in shape.text_frame.text:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None:
                        body_sizes.append(run.font.size.pt)
    assert body_sizes and min(body_sizes) >= 19


def test_presentation_write_without_target_still_multi_slide(tmp_path: Path) -> None:
    outcome = presentation_write(
        prompt="kuantum interneti sunumu",
        output_path=str(tmp_path / "serbest.pptx"),
        source_context=_RESEARCH_SUMMARY,
        overwrite=True,
    )
    assert "PPTX oluşturuldu" in outcome["text"]

    from pptx import Presentation

    deck = Presentation(str(tmp_path / "serbest.pptx"))
    # Hedef yoksa içerik kadar (≤4) + başlık + kaynaklar — en az 3 anlamlı slayt.
    assert len(list(deck.slides)) >= 3
