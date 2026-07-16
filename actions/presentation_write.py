from __future__ import annotations

import os
import re
import tempfile
import uuid
from pathlib import Path
from typing import Any

from actions._read_only_common import summarize_text
from actions._visual_block_common import (
    normalize_visual_blocks,
    render_chart_block_to_png,
    split_text_blocks,
    table_ensure_width,
    validate_visual_block_paths,
)
from actions._write_common import artifact_payload, ensure_allowed_output_path, normalize_source_context
from runtime.capability_registry import SafeCapabilityError


def _workspace_root() -> Path:
    from actions._read_only_common import workspace_root

    return workspace_root()


def _chart_temp_path() -> Path:
    directory = Path(tempfile.gettempdir())
    directory.mkdir(parents=True, exist_ok=True)
    fd, raw_path = tempfile.mkstemp(prefix=f"elyan-slide-chart-{uuid.uuid4().hex[:8]}-", suffix=".png", dir=str(directory))
    os.close(fd)
    return Path(raw_path)


def _normalize_slides(slides: Any) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    if not isinstance(slides, list):
        return normalized
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        normalized.append(
            {
                "title": str(slide.get("title", "") or "").strip(),
                "bullets": [str(item) for item in (slide.get("bullets") or []) if str(item).strip()],
                "body": str(slide.get("body", "") or "").strip(),
                "blocks": slide.get("blocks") if isinstance(slide.get("blocks"), list) else [],
            }
        )
    return normalized


def _slide_blocks(
    slide: dict[str, Any],
    *,
    selected_paths: list[str] | None = None,
) -> list[dict[str, Any]]:
    merged: list[Any] = []
    merged.extend(slide.get("blocks", []) or [])
    body = str(slide.get("body", "") or "").strip()
    if body:
        merged.extend(split_text_blocks(body))
    bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
    if bullets:
        merged.extend(split_text_blocks("\n".join(f"• {bullet}" for bullet in bullets)))
    return validate_visual_block_paths(
        normalize_visual_blocks(merged, fallback_text=""),
        selected_paths=selected_paths,
        root_resolver=_workspace_root,
    )


def _add_slide_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    bold: bool = False,
    font_size: int | None = None,
) -> None:
    from pptx.enum.text import PP_ALIGN  # type: ignore[reportMissingImports]
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size or (28 if bold else 18))
    run.font.bold = bold
    p.alignment = PP_ALIGN.LEFT


def _add_title_slide(slide: Any, title: str, subtitle: str, *, slide_width: int, slide_height: int) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[reportMissingImports]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[reportMissingImports]
    from pptx.util import Inches

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75),
        Inches(1.15),
        Inches(0.12),
        Inches(3.75),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(31, 122, 118)
    accent.line.fill.background()
    _add_slide_textbox(
        slide,
        Inches(1.15),
        Inches(1.35),
        slide_width - Inches(2.0),
        Inches(1.6),
        title,
        bold=True,
        font_size=34,
    )
    if subtitle.strip() and subtitle.strip().casefold() != title.strip().casefold():
        _add_slide_textbox(
            slide,
            Inches(1.18),
            Inches(3.15),
            slide_width - Inches(2.2),
            Inches(1.35),
            summarize_text(subtitle, max_chars=180),
            font_size=19,
        )


def _add_bullet_slide(
    slide: Any,
    title: str,
    bullets: list[str],
    *,
    slide_width: int,
    slide_height: int,
    page_number: int,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[reportMissingImports]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[reportMissingImports]
    from pptx.util import Inches, Pt

    content_left = Inches(0.85)
    content_width = slide_width - Inches(1.7)
    _add_slide_textbox(
        slide,
        content_left,
        Inches(0.48),
        content_width,
        Inches(0.72),
        title,
        bold=True,
        font_size=28,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        content_left,
        Inches(1.22),
        Inches(1.1),
        Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(31, 122, 118)
    accent.line.fill.background()

    cleaned = [summarize_text(item, max_chars=175) for item in bullets if str(item).strip()][:5]
    textbox = slide.shapes.add_textbox(
        content_left,
        Inches(1.58),
        content_width,
        slide_height - Inches(2.25),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    font_size = 21 if len(cleaned) <= 4 else 19
    for index, item in enumerate(cleaned):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(14)
        paragraph.line_spacing = 1.08
        if paragraph.runs:
            paragraph.runs[0].font.size = Pt(font_size)
            paragraph.runs[0].font.color.rgb = RGBColor(36, 42, 48)

    _add_slide_textbox(
        slide,
        slide_width - Inches(0.8),
        slide_height - Inches(0.45),
        Inches(0.35),
        Inches(0.22),
        str(page_number),
        font_size=10,
    )


def _add_slide_blocks(
    slide: Any,
    blocks: list[dict[str, Any]],
    title: str,
    *,
    slide_width: int,
    temp_paths: list[Path] | None = None,
) -> None:
    from pptx.util import Inches, Pt

    content_left = Inches(0.6)
    content_top = Inches(1.2)
    content_width = slide_width - Inches(1.2)
    current_top = content_top

    if title.strip():
        _add_slide_textbox(
            slide,
            content_left,
            Inches(0.35),
            content_width,
            Inches(0.6),
            title.strip(),
            bold=True,
            font_size=28,
        )

    for block in blocks:
        kind = str(block.get("kind", "") or "").strip().lower()
        if kind == "spacer":
            current_top += Inches(max(0.08, min(0.7, float(block.get("height", 16) or 16) / 160.0)))
            continue
        if kind == "text":
            text = str(block.get("text", "") or "").strip()
            if not text:
                continue
            level = int(block.get("level", 0) or 0)
            height = Inches(0.55 if level > 0 else 0.9)
            textbox = slide.shapes.add_textbox(content_left, current_top, content_width, height)
            frame = textbox.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(20 if level > 0 else 18)
            run.font.bold = level > 0
            current_top += height + Inches(0.08)
            continue
        if kind == "image":
            path = Path(str(block.get("path", "") or "")).expanduser()
            if not path.exists():
                continue
            try:
                from PIL import Image

                with Image.open(path) as handle:
                    img_w, img_h = handle.size
                max_width = float(content_width)
                target_width = max_width
                target_height = target_width * (img_h / max(1, img_w))
                max_height = float(Inches(3.6))
                if target_height > max_height:
                    target_height = max_height
                    target_width = target_height * (img_w / max(1, img_h))
            except Exception:
                target_width = float(content_width)
                target_height = float(Inches(2.4))
            slide.shapes.add_picture(str(path), content_left, current_top, width=int(target_width), height=int(target_height))
            current_top += int(target_height) + int(Inches(0.12))
            continue
        if kind == "chart":
            chart_path = _chart_temp_path()
            try:
                render_chart_block_to_png(block, chart_path)
                if temp_paths is not None:
                    temp_paths.append(chart_path)
            except Exception:
                try:
                    chart_path.unlink(missing_ok=True)
                except Exception:
                    pass
                continue
            try:
                from PIL import Image

                with Image.open(chart_path) as handle:
                    img_w, img_h = handle.size
                max_width = float(content_width)
                target_width = max_width
                target_height = target_width * (img_h / max(1, img_w))
                max_height = float(Inches(3.6))
                if target_height > max_height:
                    target_height = max_height
                    target_width = target_height * (img_w / max(1, img_h))
            except Exception:
                target_width = float(content_width)
                target_height = float(Inches(2.4))
            slide.shapes.add_picture(str(chart_path), content_left, current_top, width=int(target_width), height=int(target_height))
            current_top += int(target_height) + int(Inches(0.12))
            continue
        if kind == "table":
            headers, rows = table_ensure_width(
                [str(item) for item in (block.get("headers", []) or [])],
                [[str(cell) for cell in row] for row in (block.get("rows", []) or [])],
            )
            if not rows and not headers:
                continue
            table_title = str(block.get("title", "") or "").strip()
            if table_title:
                _add_slide_textbox(slide, content_left, current_top, content_width, Inches(0.35), table_title, bold=True)
                current_top += Inches(0.38)
            row_count = len(rows) + 1
            col_count = max(1, len(headers))
            table_height = Inches(min(4.2, 0.55 + row_count * 0.42))
            table_shape = slide.shapes.add_table(row_count, col_count, content_left, current_top, content_width, table_height)
            table = table_shape.table
            for index, header in enumerate(headers):
                cell = table.cell(0, index)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(12)
            for row_index, row in enumerate(rows, start=1):
                for col_index, cell_value in enumerate(row):
                    table.cell(row_index, col_index).text = cell_value
            current_top += table_height + Inches(0.15)
            continue


_SLIDE_COUNT_RE = re.compile(r"(\d+)\s*(?:sayfa|slayt|slide|page)", re.IGNORECASE)


def _requested_slide_count(*texts: str) -> int:
    """Komuttan istenen slayt/sayfa hedefi ("5 sayfalık sunum") — yoksa 0."""
    for text in texts:
        match = _SLIDE_COUNT_RE.search(str(text or ""))
        if match:
            return max(2, min(12, int(match.group(1))))
    return 0


def _derive_slide_specs(seed_text: str, deck_title: str, target: int) -> list[dict[str, Any]]:
    """`slides` verilmediğinde içerikten çok-slaytlı deste türetir.

    Eskiden başlık + TEK içerik slaytı basılıyordu; "5 sayfalık sunum" isteği
    2 slaytlık, içeriği kopyalanmış bir dosya üretiyordu. Araştırma özeti
    numaralı bulgular + "Kaynaklar:" kuyruğu biçiminde gelir; bulgular istenen
    sayıda slayta sırayla dağıtılır, kaynaklar ayrı kapanış slaytı olur.
    """
    text = str(seed_text or "").strip()
    if not text:
        return []

    sources_text = ""
    main = text
    marker = re.search(r"\n?\s*Kaynaklar\s*:\s*", main, flags=re.IGNORECASE)
    if marker:
        sources_text = main[marker.end():].strip()
        main = main[: marker.start()].strip()

    items = [part.strip(" -•\t\n") for part in re.split(r"\n\s*(?=\d+[.)]\s)", main) if part.strip()]
    if len(items) <= 1:
        items = [part.strip() for part in re.split(r"(?<=[.!?])\s+", main) if len(part.strip()) > 3]
    deduplicated: list[str] = []
    seen: set[str] = set()
    for item in items:
        normalized = " ".join(re.sub(r"^\d+[.)]?\s*", "", item).split()).strip()
        key = normalized.casefold()
        if normalized and key not in seen:
            seen.add(key)
            deduplicated.append(normalized)
    items = deduplicated
    if not items:
        return []

    has_sources = bool(sources_text)
    if target:
        content_slots = max(1, target - 1 - (1 if has_sources else 0))
    else:
        content_slots = min(4, len(items))
    content_slots = min(content_slots, len(items))

    max_items = content_slots * 4
    if len(items) > max_items:
        indexes = {
            round(index * (len(items) - 1) / max(1, max_items - 1))
            for index in range(max_items)
        }
        items = [item for index, item in enumerate(items) if index in indexes]

    chunks: list[list[str]] = [[] for _ in range(content_slots)]
    for index, item in enumerate(items):
        slot = min(index * content_slots // len(items), content_slots - 1)
        chunks[slot].append(item)

    specs: list[dict[str, Any]] = []
    for chunk in chunks:
        if not chunk:
            continue
        lead = chunk[0]
        slide_title = re.split(r"[:.;,]", lead, maxsplit=1)[0].strip()
        title_words = slide_title.split()
        if len(title_words) > 8:
            slide_title = " ".join(title_words[:8])
        if len(slide_title) > 54 or not slide_title:
            slide_title = summarize_text(lead, max_chars=54).rstrip("…")
        bullets = [summarize_text(entry, max_chars=175) for entry in chunk[:4]]
        specs.append({"title": slide_title or deck_title, "bullets": bullets, "body": ""})

    if has_sources:
        source_lines = [line.strip(" -•\t") for line in re.split(r"\n|(?:\s-\s)", sources_text) if line.strip(" -•\t")]
        specs.append({
            "title": "Kaynaklar",
            "bullets": [summarize_text(line, max_chars=180) for line in source_lines[:8]],
            "body": "",
        })
    return specs


def _single_slide_payload(
    title: str,
    blocks: list[dict[str, Any]],
    summary_text: str,
) -> list[dict[str, Any]]:
    if blocks:
        return [{"title": title, "blocks": blocks}]
    return [{"title": title, "body": summary_text, "bullets": []}]


def presentation_write(
    prompt: str = "",
    output_path: str = "",
    title: str = "",
    slides: list[dict[str, Any]] | None = None,
    blocks: list[dict[str, Any]] | None = None,
    source_context: str = "",
    overwrite: bool = False,
    _selectedPaths: list[str] | None = None,
) -> dict[str, Any]:
    from pptx import Presentation  # type: ignore[reportMissingImports]

    resolved_output = ensure_allowed_output_path(
        output_path,
        extension=".pptx",
        overwrite=overwrite,
        hint=title or prompt or "elyan-presentation",
        root_resolver=_workspace_root,
    )
    seed_text = str(source_context or prompt or "").strip()
    normalized_slides = _normalize_slides(slides)
    normalized_blocks = validate_visual_block_paths(
        normalize_visual_blocks(blocks or [], fallback_text=seed_text),
        selected_paths=_selectedPaths,
        root_resolver=_workspace_root,
    )
    if not seed_text and not normalized_slides and not normalized_blocks:
        raise SafeCapabilityError("INVALID_ARGUMENT", "Sunum oluşturmak için içerik veya slayt verisi gerekli.")

    presentation = Presentation()
    from pptx.util import Inches

    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    deck_title = str(title or resolved_output.stem).strip() or "Elyan Presentation"
    blank_layout = presentation.slide_layouts[6]
    title_slide = presentation.slides.add_slide(blank_layout)
    temp_paths: list[Path] = []
    _add_title_slide(
        title_slide,
        deck_title,
        summarize_text(seed_text or deck_title, max_chars=180),
        slide_width=int(presentation.slide_width),
        slide_height=int(presentation.slide_height),
    )

    if not normalized_slides:
        derived = _derive_slide_specs(seed_text, deck_title, _requested_slide_count(prompt, title))
        if derived:
            normalized_slides = _normalize_slides(derived)
    slide_specs = normalized_slides if normalized_slides else _single_slide_payload(deck_title, normalized_blocks, seed_text)
    try:
        for slide in slide_specs:
            page = presentation.slides.add_slide(blank_layout)
            slide_bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
            if slide_bullets and not slide.get("blocks") and not str(slide.get("body", "") or "").strip():
                _add_bullet_slide(
                    page,
                    str(slide.get("title", "") or deck_title),
                    slide_bullets,
                    slide_width=int(presentation.slide_width),
                    slide_height=int(presentation.slide_height),
                    page_number=len(presentation.slides),
                )
                continue
            page_blocks = _slide_blocks(slide, selected_paths=_selectedPaths)
            if not page_blocks and str(slide.get("body", "") or "").strip():
                page_blocks = normalize_visual_blocks([slide.get("body", "")], fallback_text=slide.get("body", ""))
            _add_slide_blocks(
                page,
                page_blocks,
                str(slide.get("title", "") or deck_title),
                slide_width=int(presentation.slide_width),
                temp_paths=temp_paths,
            )

        presentation.save(str(resolved_output))
    finally:
        for path in temp_paths:
            try:
                path.unlink(missing_ok=True)
            except Exception:
                continue
    return {
        "text": f"PPTX oluşturuldu: {resolved_output.name}",
        "result": {
            "kind": "presentation_write",
            "sourceContext": normalize_source_context(seed_text or deck_title),
            "outputPath": str(resolved_output),
            "contentType": artifact_payload(resolved_output)["contentType"],
            "created": True,
            "title": deck_title,
            "slideCount": len(presentation.slides),
            "summary": summarize_text(seed_text or deck_title, max_chars=260),
        },
        "artifacts": [artifact_payload(resolved_output)],
    }
