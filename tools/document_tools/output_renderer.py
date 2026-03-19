"""Unified document rendering for sectioned content."""

from __future__ import annotations

import html
import io
import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from tools.office_tools.content_manifest import (
    manifest_to_excel_payload,
    manifest_to_presentation_sections,
)
from utils.logger import get_logger

logger = get_logger("document.output_renderer")


@dataclass
class DocumentParagraph:
    text: str
    claim_ids: list[str] = field(default_factory=list)


@dataclass
class DocumentSection:
    title: str
    paragraphs: list[DocumentParagraph] = field(default_factory=list)


@dataclass
class SectionedDocument:
    title: str
    sections: list[DocumentSection] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_plain_text(self) -> str:
        lines = [self.title]
        for section in list(self.sections or []):
            heading = str(section.title or "").strip()
            if heading:
                lines.extend(["", heading])
            for paragraph in list(section.paragraphs or []):
                text = str(paragraph.text or "").strip()
                if text:
                    lines.extend(["", text])
        return "\n".join(lines).strip()


class UnsupportedFormat(ValueError):
    pass


class BaseRenderer:
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        raise NotImplementedError


class MarkdownRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        target = Path(path).expanduser().resolve()
        lines = [f"# {document.title}"]
        for section in list(document.sections or []):
            heading = str(section.title or "").strip()
            if heading:
                lines.extend(["", f"## {heading}", ""])
            for paragraph in list(section.paragraphs or []):
                text = str(paragraph.text or "").strip()
                if text:
                    lines.append(text)
                    lines.append("")
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")
        return {"success": True, "path": str(target), "format": "md"}


class HtmlRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        target = Path(path).expanduser().resolve()
        body = [f"<h1>{html.escape(document.title)}</h1>"]
        for section in list(document.sections or []):
            heading = str(section.title or "").strip()
            body.append("<section>")
            if heading:
                body.append(f"<h2>{html.escape(heading)}</h2>")
            for paragraph in list(section.paragraphs or []):
                body.append(f"<p>{html.escape(str(paragraph.text or '').strip())}</p>")
            body.append("</section>")
        payload = (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
            f"<title>{html.escape(document.title)}</title>"
            "<style>"
            "body{font-family:Georgia,'Times New Roman',serif;background:#ffffff;color:#111827;max-width:760px;margin:40px auto;padding:0 24px;line-height:1.65;}"
            "h1{font-size:30px;margin:0 0 24px;color:#0f172a;}"
            "h2{font-size:18px;margin:28px 0 10px;color:#1e3a8a;}"
            "p{margin:0 0 14px;}"
            "</style></head><body>"
            + "".join(body)
            + "</body></html>"
        )
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(payload, encoding="utf-8")
        return {"success": True, "path": str(target), "format": "html"}


class DocxRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        from tools.office_tools.word_tools import write_word

        paragraphs: list[str] = []
        for section in list(document.sections or []):
            paragraphs.append(str(section.title or "").strip())
            for paragraph in list(section.paragraphs or []):
                text = str(paragraph.text or "").strip()
                if text:
                    paragraphs.append(text)
        return await write_word(
            path=str(path),
            title=document.title,
            content=document.to_plain_text(),
            paragraphs=paragraphs,
        )


class XlsxRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        from tools.office_tools.excel_tools import write_excel

        manifest = document.metadata.get("office_content_manifest") if isinstance(document.metadata, dict) else None
        if isinstance(manifest, dict):
            data, headers = manifest_to_excel_payload(manifest)
            if data:
                return await write_excel(
                    path=str(path),
                    data=data,
                    headers=headers,
                    multi_sheet=True,
                )

        rows: list[dict[str, Any]] = [{"Content": document.title}]
        for section in list(document.sections or []):
            heading = str(section.title or "").strip()
            if heading:
                rows.append({"Content": heading})
            for paragraph in list(section.paragraphs or []):
                text = str(paragraph.text or "").strip()
                if text:
                    rows.append({"Content": text})
        return await write_excel(path=str(path), data=rows, headers=["Content"], sheet_name="Content")


class PdfRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        try:
            from reportlab.lib.colors import HexColor
            from reportlab.lib.enums import TA_LEFT
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
        except Exception as exc:
            return {"success": False, "error": f"reportlab unavailable: {exc}"}

        target = Path(path).expanduser().resolve()
        target.parent.mkdir(parents=True, exist_ok=True)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "elyan_title",
            parent=styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=19,
            leading=24,
            textColor=HexColor("#0f172a"),
            alignment=TA_LEFT,
            spaceAfter=20,
        )
        heading_style = ParagraphStyle(
            "elyan_heading",
            parent=styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=17,
            textColor=HexColor("#1e3a8a"),
            spaceBefore=12,
            spaceAfter=8,
        )
        body_style = ParagraphStyle(
            "elyan_body",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.5,
            leading=15,
            textColor=HexColor("#111827"),
            spaceAfter=10,
        )
        story: list[Any] = [Paragraph(html.escape(document.title), title_style), Spacer(1, 8)]
        for section in list(document.sections or []):
            heading = str(section.title or "").strip()
            if heading:
                story.append(Paragraph(html.escape(heading), heading_style))
            for paragraph in list(section.paragraphs or []):
                text = str(paragraph.text or "").strip()
                if text:
                    story.append(Paragraph(html.escape(text), body_style))
        if len(story) <= 2:
            story.append(Paragraph(html.escape(document.to_plain_text()), body_style))
        doc = SimpleDocTemplate(str(target), pagesize=A4, leftMargin=48, rightMargin=48, topMargin=54, bottomMargin=54)
        doc.title = document.title
        doc.build(story)
        return {"success": True, "path": str(target), "format": "pdf"}


class PptxRenderer(BaseRenderer):
    async def render_to_path(self, document: SectionedDocument, path: str) -> dict[str, Any]:
        try:
            from pptx import Presentation
        except Exception as exc:
            return {"success": False, "error": f"python-pptx unavailable: {exc}"}

        target = Path(path).expanduser().resolve()
        target.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        manifest = document.metadata.get("office_content_manifest") if isinstance(document.metadata, dict) else None
        slide_sections = list(document.sections or [])
        subtitle_text = "Generated by Elyan"
        if isinstance(manifest, dict):
            slide_sections = manifest_to_presentation_sections(manifest)
            summary = manifest.get("summary") if isinstance(manifest.get("summary"), dict) else {}
            subtitle_text = str(
                summary.get("preview")
                or summary.get("brief")
                or summary.get("quality_status")
                or manifest.get("brief")
                or "Generated by Elyan"
            ).strip()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = document.title
        subtitle = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
        if subtitle is not None:
            subtitle.text = subtitle_text[:240] if subtitle_text else "Generated by Elyan"
        for section in slide_sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(section.title or "").strip()
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, paragraph in enumerate(list(section.paragraphs or [])):
                p = body.paragraphs[0] if index == 0 else body.add_paragraph()
                p.text = str(paragraph.text or "").strip()
        prs.save(str(target))
        return {"success": True, "path": str(target), "format": "pptx"}


class DocumentRenderer:
    RENDERERS = {
        "docx": DocxRenderer,
        "pdf": PdfRenderer,
        "xlsx": XlsxRenderer,
        "md": MarkdownRenderer,
        "html": HtmlRenderer,
        "pptx": PptxRenderer,
    }

    def get_renderer(self, fmt: str) -> BaseRenderer:
        clean = str(fmt or "").strip().lower()
        renderer_cls = self.RENDERERS.get(clean)
        if renderer_cls is None:
            raise UnsupportedFormat(clean)
        return renderer_cls()

    async def render_to_path(self, document: SectionedDocument, fmt: str, path: str) -> dict[str, Any]:
        renderer = self.get_renderer(fmt)
        return await renderer.render_to_path(document, path)


def sections_to_sectioned_document(
    *,
    title: str,
    sections: list[dict[str, Any]],
    metadata: dict[str, Any] | None = None,
) -> SectionedDocument:
    rows: list[DocumentSection] = []
    for section in list(sections or []):
        if not isinstance(section, dict):
            continue
        title_text = str(section.get("title") or "").strip()
        paragraphs = []
        for paragraph in list(section.get("paragraphs") or []):
            if isinstance(paragraph, dict):
                paragraphs.append(
                    DocumentParagraph(
                        text=str(paragraph.get("text") or "").strip(),
                        claim_ids=[str(item).strip() for item in list(paragraph.get("claim_ids") or []) if str(item).strip()],
                    )
                )
            elif str(paragraph or "").strip():
                paragraphs.append(DocumentParagraph(text=str(paragraph).strip()))
        rows.append(DocumentSection(title=title_text, paragraphs=paragraphs))
    return SectionedDocument(title=str(title or "").strip(), sections=rows, metadata=dict(metadata or {}))


def sectioned_document_json(document: SectionedDocument) -> str:
    payload = {
        "title": document.title,
        "metadata": dict(document.metadata or {}),
        "sections": [
            {
                "title": section.title,
                "paragraphs": [{"text": paragraph.text, "claim_ids": list(paragraph.claim_ids or [])} for paragraph in section.paragraphs],
            }
            for section in list(document.sections or [])
        ],
    }
    return json.dumps(payload, ensure_ascii=False, indent=2)


__all__ = [
    "DocumentParagraph",
    "DocumentRenderer",
    "DocumentSection",
    "SectionedDocument",
    "UnsupportedFormat",
    "sectioned_document_json",
    "sections_to_sectioned_document",
]
